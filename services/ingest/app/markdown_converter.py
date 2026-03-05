"""
============================================================================
FILE: services/ingest/app/markdown_converter.py
PURPOSE: ★ SIGNATURE OPTIMIZATION — Convert ALL document formats to Markdown
         BEFORE chunking and embedding. This is the single most impactful
         accuracy optimization in the entire RAG pipeline.
ARCHITECTURE REF: §3.1 — Convert-to-Markdown Before Embedding
DEPENDENCIES: pymupdf, python-docx, openpyxl, python-pptx
============================================================================

WHY Convert to Markdown First?
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
1. ACCURACY BOOST: Markdown preserves document structure (headings, lists,
   tables) in a format that BGE-M3 understands natively. Raw text extraction
   from PDFs loses hierarchical context.

2. CONSISTENCY: All documents — regardless of source format — go through the
   same downstream pipeline (chunker.py → embedding_service). No format-specific
   edge cases in later stages.

3. METADATA PRESERVATION: Section headings from markdown structure become
   metadata tags on chunks, enabling better citation attribution:
   "Source: HR_Policy.pdf | Section: Annual Leave | Page: 5"

4. DEBUGGABILITY: Stored in MinIO as <original_name>.md — admins can inspect
   the converted markdown to verify extraction quality before investing in
   re-processing.

Output Format:
━━━━━━━━━━━━━
Every converted document starts with YAML frontmatter:
    ---
    filename: original_document.pdf
    format: pdf
    page_count: 12
    upload_date: 2026-02-20
    ---

Followed by document content with:
- # ## ### headings preserving document hierarchy
- | pipe | table | formatting |
- - bullet list items
- **bold** and *italic* text
- <!-- PAGE_BREAK: page_N --> markers for page-level citation
"""

import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def _build_frontmatter(filename: str, fmt: str, page_count: Optional[int] = None) -> str:
    """
    Build YAML frontmatter block for the markdown document.

    The frontmatter is parsed by metadata_extractor.py to tag each chunk
    with document-level metadata (filename, upload date, etc.).

    Args:
        filename: Original filename of the uploaded document.
        fmt: File format (pdf, docx, xlsx, pptx, txt, md).
        page_count: Number of pages/slides/sheets (None if not applicable).

    Returns:
        YAML frontmatter string including the closing '---' delimiter.
    """
    today = datetime.utcnow().strftime("%Y-%m-%d")
    lines = [
        "---",
        f"filename: {filename}",
        f"format: {fmt}",
        f"upload_date: {today}",
    ]
    if page_count is not None:
        lines.append(f"page_count: {page_count}")
    lines.append("---")
    lines.append("")  # blank line after frontmatter
    return "\n".join(lines)


def pdf_to_markdown(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """
    Convert a PDF to structured Markdown using PyMuPDF.

    Strategy:
    1. PyMuPDF (fitz) extracts text blocks with position and font information
    2. Font size analysis maps font sizes to heading levels (# ## ###)
    3. Tables extracted using PyMuPDF's built-in table detector
    4. Page break markers inserted between pages for citation tracking

    Args:
        file_bytes: Raw PDF file content.
        filename: Original filename (for frontmatter).

    Returns:
        Tuple of (markdown_string, page_count).

    Raises:
        ValueError: If the PDF cannot be parsed.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("pymupdf is not installed. Add it to requirements.txt.")

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        raise ValueError(f"Failed to open PDF: {exc}") from exc

    page_count = len(doc)
    md_sections = [_build_frontmatter(filename, "pdf", page_count)]

    for page_num, page in enumerate(doc, start=1):
        # Page break marker — enables page-level source citations
        if page_num > 1:
            md_sections.append(f"\n<!-- PAGE_BREAK: page_{page_num} -->\n")

        # Extract text with block structure preserved
        # "dict" mode returns blocks with position, font, and text data
        page_dict = page.get_text("dict")

        # Determine the dominant font size on this page (for relative heading detection)
        font_sizes = []
        for block in page_dict.get("blocks", []):
            if block["type"] != 0:  # 0 = text block
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    if span.get("text", "").strip():
                        font_sizes.append(span["size"])

        # Median font size = body text; larger = headings
        if font_sizes:
            font_sizes.sort()
            median_size = font_sizes[len(font_sizes) // 2]
        else:
            median_size = 12.0  # Default assumption

        # Extract tables from this page
        tables = page.find_tables()
        table_bboxes = [t.bbox for t in tables.tables] if tables.tables else []

        # Process each text block
        page_text_parts = []
        for block in page_dict.get("blocks", []):
            if block["type"] != 0:  # Skip image blocks
                continue

            # Check if this block overlaps with a table (skip — tables handled separately)
            block_rect = fitz.Rect(block["bbox"])
            in_table = any(block_rect.intersects(fitz.Rect(tb)) for tb in table_bboxes)
            if in_table:
                continue

            block_text_parts = []
            max_font_in_block = 0.0

            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if not text:
                        continue
                    font_size = span.get("size", 12.0)
                    max_font_in_block = max(max_font_in_block, font_size)

                    # Preserve bold/italic formatting
                    flags = span.get("flags", 0)
                    is_bold   = bool(flags & 16)  # Bold flag in PyMuPDF
                    is_italic = bool(flags & 2)   # Italic flag

                    if is_bold and is_italic:
                        text = f"***{text}***"
                    elif is_bold:
                        text = f"**{text}**"
                    elif is_italic:
                        text = f"*{text}*"

                    block_text_parts.append(text)

            if not block_text_parts:
                continue

            full_text = " ".join(block_text_parts)

            # Map font size to heading level
            # Threshold: > 1.5× median = H1, > 1.3× = H2, > 1.1× = H3, else body
            if max_font_in_block > median_size * 1.5:
                page_text_parts.append(f"# {full_text}")
            elif max_font_in_block > median_size * 1.3:
                page_text_parts.append(f"## {full_text}")
            elif max_font_in_block > median_size * 1.1:
                page_text_parts.append(f"### {full_text}")
            else:
                page_text_parts.append(full_text)

        # Append extracted tables in markdown format
        for table in (tables.tables if tables.tables else []):
            table_md = _pymupdf_table_to_markdown(table)
            if table_md:
                page_text_parts.append(table_md)

        md_sections.append("\n\n".join(page_text_parts))

    doc.close()
    return "\n\n".join(md_sections), page_count


def _pymupdf_table_to_markdown(table) -> str:
    """
    Convert a PyMuPDF table object to a Markdown pipe table.

    Args:
        table: PyMuPDF Table object from page.find_tables().

    Returns:
        Markdown table string, or empty string if table is empty.
    """
    try:
        rows = table.extract()
        if not rows:
            return ""

        md_rows = []
        for row_idx, row in enumerate(rows):
            # Clean cell text (remove None, strip whitespace)
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            md_rows.append("| " + " | ".join(cells) + " |")

            # Add header separator after first row
            if row_idx == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                md_rows.append(separator)

        return "\n".join(md_rows)
    except Exception as exc:
        logger.warning(f"Failed to convert table to markdown: {exc}")
        return ""


def docx_to_markdown(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """
    Convert a DOCX file to Markdown preserving heading levels, lists, and tables.

    Strategy:
    - python-docx provides structured access to paragraphs and their styles
    - Heading styles (Heading 1, Heading 2, etc.) → # ## ### headings
    - List paragraphs → - bullet items or 1. numbered items
    - Tables → | pipe | tables |
    - Runs with bold/italic → **bold** / *italic*

    Args:
        file_bytes: Raw DOCX file content.
        filename: Original filename (for frontmatter).

    Returns:
        Tuple of (markdown_string, paragraph_count_as_proxy_for_pages).
    """
    try:
        import io
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        raise ImportError("python-docx is not installed.")

    doc = Document(io.BytesIO(file_bytes))
    md_parts = [_build_frontmatter(filename, "docx")]

    para_count = 0
    for element in doc.element.body:
        # Paragraphs
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            from docx.oxml.ns import qn
            from docx import Document
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name if para.style else ""
            para_count += 1

            # Map DOCX heading styles to Markdown heading levels
            if style_name.startswith("Heading 1"):
                md_parts.append(f"# {text}")
            elif style_name.startswith("Heading 2"):
                md_parts.append(f"## {text}")
            elif style_name.startswith("Heading 3"):
                md_parts.append(f"### {text}")
            elif style_name.startswith("Heading 4"):
                md_parts.append(f"#### {text}")
            elif "List Bullet" in style_name:
                md_parts.append(f"- {text}")
            elif "List Number" in style_name:
                md_parts.append(f"1. {text}")  # numbering is handled by markdown renderer
            else:
                # Process runs to capture bold/italic within the paragraph
                rich_text = _process_docx_runs(para)
                md_parts.append(rich_text)

        elif tag == "tbl":
            # Table — render as markdown pipe table
            from docx.table import Table
            table = Table(element, doc)
            table_md = _docx_table_to_markdown(table)
            if table_md:
                md_parts.append(table_md)

    # Estimate page count (very rough — 40 paragraphs ≈ 1 page)
    estimated_pages = max(1, para_count // 40)

    return "\n\n".join(md_parts), estimated_pages


def _process_docx_runs(para) -> str:
    """Extract text from a paragraph with inline bold/italic formatting."""
    parts = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        if run.bold and run.italic:
            parts.append(f"***{text}***")
        elif run.bold:
            parts.append(f"**{text}**")
        elif run.italic:
            parts.append(f"*{text}*")
        else:
            parts.append(text)
    return "".join(parts) or para.text


def _docx_table_to_markdown(table) -> str:
    """Convert a python-docx Table to a Markdown pipe table."""
    if not table.rows:
        return ""

    md_rows = []
    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        md_rows.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")

    return "\n".join(md_rows)


def xlsx_to_markdown(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """
    Convert an Excel XLSX file to Markdown.

    Each sheet becomes a level-2 heading section.
    Sheet data is rendered as a Markdown pipe table.
    Empty rows and columns are stripped for cleanliness.

    Args:
        file_bytes: Raw XLSX file content.
        filename: Original filename (for frontmatter).

    Returns:
        Tuple of (markdown_string, sheet_count).
    """
    try:
        import io
        import openpyxl
    except ImportError:
        raise ImportError("openpyxl is not installed.")

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheet_count = len(wb.sheetnames)
    md_parts = [_build_frontmatter(filename, "xlsx", sheet_count)]

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md_parts.append(f"## Sheet: {sheet_name}")

        # Read all rows, filtering out empty rows
        rows = []
        for row in ws.iter_rows(values_only=True):
            # Skip fully empty rows
            if all(cell is None or str(cell).strip() == "" for cell in row):
                continue
            rows.append([str(cell).strip() if cell is not None else "" for cell in row])

        if not rows:
            md_parts.append("*(Empty sheet)*")
            continue

        # Remove fully-empty trailing columns
        max_col = 0
        for row in rows:
            for col_idx, cell in enumerate(row):
                if cell:
                    max_col = max(max_col, col_idx)
        rows = [row[:max_col + 1] for row in rows]

        # Build markdown table
        md_rows = []
        for row_idx, row in enumerate(rows):
            md_rows.append("| " + " | ".join(row) + " |")
            if row_idx == 0:
                md_rows.append("| " + " | ".join(["---"] * len(row)) + " |")

        md_parts.append("\n".join(md_rows))

    wb.close()
    return "\n\n".join(md_parts), sheet_count


def pptx_to_markdown(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """
    Convert a PowerPoint PPTX file to Markdown.

    Each slide becomes a level-2 heading section.
    Bullet points from text frames are extracted.
    Speaker notes are included as blockquotes.

    Args:
        file_bytes: Raw PPTX file content.
        filename: Original filename (for frontmatter).

    Returns:
        Tuple of (markdown_string, slide_count).
    """
    try:
        import io
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("python-pptx is not installed.")

    prs = Presentation(io.BytesIO(file_bytes))
    slide_count = len(prs.slides)
    md_parts = [_build_frontmatter(filename, "pptx", slide_count)]

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Get slide title (if present)
        title_text = ""
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()

        slide_heading = f"## Slide {slide_num}"
        if title_text:
            slide_heading += f": {title_text}"
        md_parts.append(slide_heading)

        # Extract text from all non-title text frames
        slide_content = []
        for shape in slide.shapes:
            # Skip the title shape (already handled above)
            if shape == slide.shapes.title:
                continue

            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detect bullet level from paragraph indent level
                level = para.level  # 0 = top level, 1+ = nested
                indent = "  " * level
                slide_content.append(f"{indent}- {text}")

        if slide_content:
            md_parts.append("\n".join(slide_content))

        # Include speaker notes as blockquotes (valuable context for HR training content)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                notes_text = notes_frame.text.strip()
                # Format as blockquote (each line gets > prefix)
                blockquote = "\n".join(f"> {line}" for line in notes_text.split("\n"))
                md_parts.append(f"**Speaker Notes:**\n{blockquote}")

    return "\n\n".join(md_parts), slide_count


def txt_to_markdown(file_bytes: bytes, filename: str) -> tuple[str, int]:
    """
    Add metadata frontmatter to plain text or Markdown files.

    For TXT files: wraps content with frontmatter.
    For MD files: prepends frontmatter to existing markdown.

    Args:
        file_bytes: Raw file content (UTF-8 assumed, falls back to latin-1).
        filename: Original filename (for frontmatter).

    Returns:
        Tuple of (markdown_string, estimated_page_count).
    """
    # Try UTF-8 first, then fall back to latin-1 (common in legacy HR documents)
    try:
        content = file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        content = file_bytes.decode("latin-1")

    line_count = len(content.splitlines())
    # Rough estimate: 50 lines per page
    estimated_pages = max(1, line_count // 50)

    frontmatter = _build_frontmatter(filename, Path(filename).suffix.lstrip("."), estimated_pages)
    return frontmatter + content, estimated_pages
